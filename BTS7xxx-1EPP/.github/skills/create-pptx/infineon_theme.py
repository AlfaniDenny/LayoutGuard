"""
Infineon Technologies Official Brand Theme (2023-2)

Extracted from the official Infineon PowerPoint template.
Color scheme: "Infineon 2023 - 2"
Slide size: 13.333" × 7.500" (widescreen)

Usage:
    from infineon_theme import INFINEON_THEME, create_title_slide, create_content_slide
"""

import os
import sys

# Add bundled packages to path as fallback
_packages_dir = os.path.join(os.path.dirname(__file__), "packages")
if os.path.isdir(_packages_dir) and _packages_dir not in sys.path:
    sys.path.insert(0, _packages_dir)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════
# INFINEON COLOR PALETTE  (Infineon 2023 - 2 theme)
# ═══════════════════════════════════════════════════════════

INFINEON_THEME = {
    "colors": {
        # ── Theme-mapped colors (from template XML) ──
        "textPrimary":    "1D1D1D",   # dk1 / TEXT_1
        "white":          "FFFFFF",   # lt1 / BACKGROUND_1
        "infTeal":        "0A8276",   # dk2 / TEXT_2 — Infineon primary brand color
        "infWarmGrey":    "8D8786",   # lt2 / BACKGROUND_2
        "accent1":        "0A8276",   # accent1
        "accent2":        "575352",   # accent2
        "infOrange":      "F97414",   # accent3
        "infLime":        "9BBA43",   # accent4
        "infYellow":      "FCD442",   # accent5
        "infMagenta":     "9C216E",   # accent6

        # ── Convenience aliases ──
        "primary":        "0A8276",   # Infineon Teal – titles, headers, branding
        "black":          "1D1D1D",   # Primary text
        "grey":           "8D8786",   # Captions, secondary text
        "orange":         "F97414",   # Highlights, in-progress
        "lime":           "9BBA43",   # Positive, success
        "yellow":         "FCD442",   # Warnings, attention
        "magenta":        "9C216E",   # Accent
        "darkGrey":       "575352",   # Borders, secondary elements

        # ── Additional UI colors (from actual slides) ──
        "blue":           "0070C0",   # Used in some slides for boxes
        "tealLight":      "6CB4AD",   # Lighter teal variant
        "greenSuccess":   "3B9B91",   # Year2-style green
        "orangeWarm":     "FF9737",   # Warmer orange variant
        "lightGrey":      "F2F2F2",   # Card backgrounds
        "medGrey":        "D9D9D9",   # Borders, dividers
        "pinkLight":      "E3A9CB",   # Soft pink accent
        "yellowSoft":     "FBE29B",   # Soft yellow background
        "purpleTool":     "DCD5D7",   # Tool/disabled state

        # ── Status Colors ──
        "statusGreen":    "3B9B91",   # Completed / On Track
        "statusOrange":   "F97414",   # In Progress
        "statusRed":      "FF0000",   # Blocked / Critical
        "statusGrey":     "8D8786",   # Not Started
    },

    # ═══════════════════════════════════════════════════════════
    # TYPOGRAPHY
    # ═══════════════════════════════════════════════════════════
    "fonts": {
        "primary":    "Arial",
        "fallback":   "Arial Unicode MS",

        "sizes": {
            "titleSlide":     36,
            "slideTitle":     24,
            "sectionHeader":  16,
            "bodyLarge":      17,
            "bodyText":       12,
            "bodySmall":      10,
            "caption":        8,
            "captionSmall":   7,
            "metricLarge":    42,
            "metricSmall":    24,
        },
    },

    # ═══════════════════════════════════════════════════════════
    # LAYOUT DIMENSIONS (in inches for 13.333 × 7.500 slide)
    # ═══════════════════════════════════════════════════════════
    "layout": {
        "slideWidth":     13.333,
        "slideHeight":    7.5,

        "margins": {
            "standard":   0.37,
            "narrow":     0.2,
            "wide":       0.68,
        },

        "title": {
            "x": 0.37,  "y": 0.21,
            "w": 10.51, "h": 0.79,
        },

        "content": {
            "x": 0.37,  "y": 1.39,
            "w": 12.60, "h": 5.59,
        },

        "logo": {
            "x": 11.66, "y": 0.21,
            "w": 1.32,  "h": 0.58,
        },

        "footer": {
            "y":           7.10,
            "height":      0.40,
            "dateX":       0.37,
            "dateW":       1.20,
            "classX":      1.35,
            "classW":      1.06,
            "copyrightX":  4.97,
            "copyrightW":  3.39,
            "propX":       10.33,
            "propW":       1.06,
            "pageNumX":    12.55,
            "pageNumW":    0.42,
        },

        "titleSlide": {
            "imageX": 0,      "imageY": 0,
            "imageW": 13.333, "imageH": 4.09,
            "titleX": 0.68,   "titleY": 4.55,
            "titleW": 9.06,   "titleH": 1.22,
            "subtitleX": 0.68, "subtitleY": 6.17,
            "subtitleW": 9.06, "subtitleH": 0.67,
            "logoX": 11.20,   "logoW": 1.76,
            "logoY": 4.86,    "logoH": 0.77,
        },
    },
}


# ═══════════════════════════════════════════════════════════
# HELPER: Convert hex string to RGBColor
# ═══════════════════════════════════════════════════════════

def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a hex color string (e.g. '0A8276') to an RGBColor."""
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _color(name: str) -> RGBColor:
    """Shortcut: get RGBColor from INFINEON_THEME color name."""
    return hex_to_rgb(INFINEON_THEME["colors"][name])


# ═══════════════════════════════════════════════════════════
# PRESENTATION FACTORY
# ═══════════════════════════════════════════════════════════

def create_presentation(title: str = "Untitled", author: str = "Infineon Technologies") -> Presentation:
    """Create a new widescreen (13.333" × 7.500") presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = title
    prs.core_properties.author = author
    return prs


# ═══════════════════════════════════════════════════════════
# LOGO LOADER
# ═══════════════════════════════════════════════════════════

_cached_logo_path = None


def find_infineon_logo() -> str | None:
    """Search common locations for the Infineon logo. Returns path or None."""
    global _cached_logo_path
    if _cached_logo_path is not None:
        return _cached_logo_path if _cached_logo_path != "" else None

    candidates = [
        os.path.join(os.path.dirname(__file__), "resources", "infineon_logo.png"),
        os.path.join(os.path.dirname(__file__), "images", "infineon_logo.png"),
        os.path.join(os.path.dirname(__file__), "..", "images", "logos", "infineon_logo.png"),
        os.path.join(os.path.dirname(__file__), "infineon_logo.png"),
    ]
    for p in candidates:
        if os.path.isfile(p):
            _cached_logo_path = os.path.abspath(p)
            return _cached_logo_path

    _cached_logo_path = ""
    return None


# ═══════════════════════════════════════════════════════════
# SHAPE HELPERS  (low-level, used by slide builders)
# ═══════════════════════════════════════════════════════════

def add_textbox(slide, left, top, width, height, text,
                font_size=12, bold=False, color="1D1D1D",
                alignment=PP_ALIGN.LEFT, font_name="Arial",
                valign=MSO_ANCHOR.TOP, word_wrap=True):
    """Add a text box with sensible defaults."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = hex_to_rgb(color) if isinstance(color, str) else color
    p.font.name = font_name
    p.alignment = alignment
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_shape(slide, shape_type, left, top, width, height,
              fill_color=None, line_color=None, line_width=None):
    """Add a shape with optional fill and line."""
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color) if isinstance(fill_color, str) else fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = hex_to_rgb(line_color) if isinstance(line_color, str) else line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_rich_text(slide, left, top, width, height, rich_items,
                  font_size=12, font_name="Arial", color="1D1D1D",
                  alignment=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    """
    Add a text box with rich (formatted) text.

    rich_items: list of dicts with keys:
        - text (str): the text content
        - bold (bool, optional)
        - italic (bool, optional)
        - color (str hex, optional)
        - size (int pt, optional)
        - bullet (bool, optional): if True, add as bullet paragraph
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(rich_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = item.get("text", "")
        p.font.size = Pt(item.get("size", font_size))
        p.font.name = item.get("font", font_name)
        p.font.bold = item.get("bold", False)
        p.font.italic = item.get("italic", False)
        p.font.color.rgb = hex_to_rgb(item.get("color", color))
        p.alignment = alignment

        if item.get("bullet"):
            p.level = item.get("level", 0)

    return txBox


# ═══════════════════════════════════════════════════════════
# HEADER & FOOTER  (applied to every content slide)
# ═══════════════════════════════════════════════════════════

def add_infineon_header(slide, title, subtitle=None, logo_path=None):
    """Add the standard Infineon header bar to a content slide."""
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]
    F = INFINEON_THEME["fonts"]
    t = L["title"]

    add_textbox(slide, t["x"], t["y"], t["w"], t["h"],
                title, font_size=F["sizes"]["slideTitle"],
                bold=False, color=C["infTeal"],
                alignment=PP_ALIGN.LEFT,
                valign=MSO_ANCHOR.BOTTOM)

    if subtitle:
        add_textbox(slide, t["x"], t["y"] + t["h"] + 0.05,
                    t["w"], 0.35, subtitle,
                    font_size=F["sizes"]["bodySmall"],
                    color=C["grey"])

    if logo_path and os.path.isfile(logo_path):
        lg = L["logo"]
        slide.shapes.add_picture(logo_path,
                                 Inches(lg["x"]), Inches(lg["y"]),
                                 Inches(lg["w"]), Inches(lg["h"]))


def add_infineon_footer(slide, date_str=None, classification=None,
                        copyright_text=None, page_number=None):
    """Add the standard Infineon footer to a slide."""
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]
    F = INFINEON_THEME["fonts"]
    f = L["footer"]
    fs = F["sizes"]["caption"]

    from datetime import date as _date
    today = _date.today()
    date_str = date_str or today.isoformat()
    copyright_text = copyright_text or f"Copyright © Infineon Technologies AG {today.year}. All rights reserved."

    # Date
    add_textbox(slide, f["dateX"], f["y"], f["dateW"], f["height"],
                date_str, font_size=fs, color=C["black"])

    # Classification
    if classification:
        add_textbox(slide, f["classX"], f["y"], f["classW"], f["height"],
                    classification, font_size=fs, color=C["black"])

    # Copyright
    add_textbox(slide, f["copyrightX"], f["y"], f["copyrightW"], f["height"],
                copyright_text, font_size=fs, color=C["black"])

    # Page number (always shown)
    page_str = str(page_number) if page_number else ""
    add_textbox(slide, f["pageNumX"], f["y"], f["pageNumW"], f["height"],
                page_str, font_size=fs, color=C["black"],
                alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE CREATORS
# ═══════════════════════════════════════════════════════════

def create_title_slide(prs, title, subtitle=None,
                       logo_path=None, key_visual_path=None):
    """
    Create an Infineon title slide.

    Args:
        prs: Presentation object
        title: Main title text
        subtitle: Subtitle / department / date
        logo_path: Path to Infineon logo PNG
        key_visual_path: Path to key visual image (fills top area)

    Returns:
        slide object
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]
    F = INFINEON_THEME["fonts"]
    TS = L["titleSlide"]

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_to_rgb(C["white"])

    # Key visual area
    if key_visual_path and os.path.isfile(key_visual_path):
        slide.shapes.add_picture(
            key_visual_path,
            Inches(TS["imageX"]), Inches(TS["imageY"]),
            Inches(TS["imageW"]), Inches(TS["imageH"])
        )
    else:
        # Try to find a default title picture in resources/
        _resources_dir = os.path.join(os.path.dirname(__file__), "resources")
        _fallback_images = ["title-picture.png", "title-picture2.png", "title-picture3.jpg"]
        _fallback_path = None
        for _img in _fallback_images:
            _candidate = os.path.join(_resources_dir, _img)
            if os.path.isfile(_candidate):
                _fallback_path = _candidate
                break
        if _fallback_path:
            slide.shapes.add_picture(
                _fallback_path,
                Inches(TS["imageX"]), Inches(TS["imageY"]),
                Inches(TS["imageW"]), Inches(TS["imageH"])
            )
        else:
            # Teal placeholder as last resort
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      TS["imageX"], TS["imageY"],
                      TS["imageW"], TS["imageH"],
                      fill_color=C["infTeal"])

    # Logo
    if logo_path and os.path.isfile(logo_path):
        slide.shapes.add_picture(
            logo_path,
            Inches(TS["logoX"]), Inches(TS["logoY"]),
            Inches(TS["logoW"]), Inches(TS["logoH"])
        )

    # Title
    add_textbox(slide, TS["titleX"], TS["titleY"], TS["titleW"], TS["titleH"],
                title, font_size=F["sizes"]["titleSlide"], bold=True,
                color=C["infTeal"])

    # Subtitle
    if subtitle:
        add_textbox(slide, TS["subtitleX"], TS["subtitleY"],
                    TS["subtitleW"], TS["subtitleH"],
                    subtitle, font_size=14, color=C["grey"])

    return slide


def create_content_slide(prs, title, subtitle=None,
                         logo_path=None, page_number=None):
    """
    Create a standard content slide with header + footer.

    Returns the slide — caller adds content to it.
    """
    C = INFINEON_THEME["colors"]
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_to_rgb(C["white"])

    add_infineon_header(slide, title, subtitle, logo_path)
    add_infineon_footer(slide, page_number=page_number)
    return slide


def create_section_slide(prs, section_title, logo_path=None):
    """
    Create a section divider slide (teal background, white text).
    """
    C = INFINEON_THEME["colors"]
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_to_rgb(C["infTeal"])

    add_textbox(slide, 0.68, 3.0, 11.0, 1.5, section_title,
                font_size=36, bold=True, color=C["white"],
                alignment=PP_ALIGN.LEFT)

    if logo_path and os.path.isfile(logo_path):
        slide.shapes.add_picture(logo_path,
                                 Inches(11.20), Inches(0.40),
                                 Inches(1.76), Inches(0.77))

    return slide


def create_conclusion_slide(prs, title=None, bullet_points=None,
                            contact_email=None, logo_path=None):
    """
    Create a conclusion / closing slide.
    If title is None, creates a logo-only closing slide.
    """
    C = INFINEON_THEME["colors"]
    F = INFINEON_THEME["fonts"]
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_to_rgb(C["white"])

    if title:
        add_textbox(slide, 0.68, 0.5, 11.0, 0.9, title,
                    font_size=28, bold=True, color=C["infTeal"])

        if bullet_points:
            add_rich_text(slide, 0.68, 1.6, 11.0, 4.5,
                          [{"text": pt, "color": C["black"]}
                           for pt in bullet_points],
                          font_size=14)

        if contact_email:
            add_textbox(slide, 0.68, 6.5, 5, 0.4,
                        f"Contact: {contact_email}",
                        font_size=11, color=C["infTeal"])
    else:
        if logo_path and os.path.isfile(logo_path):
            slide.shapes.add_picture(logo_path,
                                     Inches(4.99), Inches(3.01),
                                     Inches(3.37), Inches(1.47))
        else:
            add_textbox(slide, 4.0, 2.5, 5.33, 2.5,
                        "[ Infineon Logo ]",
                        font_size=24, color=C["infTeal"],
                        alignment=PP_ALIGN.CENTER)

    return slide
