"""
Slide Builder Utilities

Pre-built layout functions for common Infineon slide types.
Uses the Infineon 2023-2 template (13.333" × 7.500" slides).

Layouts:
    - Status cards (project tracking)
    - Comparison tables
    - Architecture diagrams
    - Roadmap timelines
    - Two-column layouts
    - Metric / KPI dashboard
    - Image gallery
    - KPI charts (planned vs achieved)
    - Effort charts (planned vs delivered)
    - Burndown charts (stacked)
    - Cumulative trend charts (line)

Usage:
    from slide_builder import create_status_card, create_comparison_table, create_kpi_chart
"""

import os
import sys

# Add bundled packages to path as fallback
_packages_dir = os.path.join(os.path.dirname(__file__), "packages")
if os.path.isdir(_packages_dir) and _packages_dir not in sys.path:
    sys.path.insert(0, _packages_dir)

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

from infineon_theme import (
    INFINEON_THEME, hex_to_rgb, _color,
    add_textbox, add_shape, add_rich_text,
    add_infineon_header, add_infineon_footer,
)


# ═══════════════════════════════════════════════════════════
# STATUS CARD LAYOUT
# ═══════════════════════════════════════════════════════════

def create_status_card(slide, *,
                       title,
                       description,
                       status,
                       status_color,
                       status_details=None,
                       users=None,
                       metrics=None,
                       logo_path=None,
                       page_number=None):
    """
    Create a status card layout for project tracking.

    Args:
        slide: slide object (from prs.slides.add_slide)
        title (str): slide title
        description (str | list[dict]): description text or rich text items
        status (str): status label, e.g. "Under Development"
        status_color (str): hex color for the status badge
        status_details (str | list[dict], optional): extra status info
        users (list[str], optional): user / team descriptions
        metrics (list[dict], optional): each dict has keys: label, value, color
        logo_path (str, optional): path to logo PNG
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]

    add_infineon_header(slide, title, logo_path=logo_path)
    add_infineon_footer(slide, page_number=page_number)

    # ─── Left Panel: Card ───
    cardX, cardY, cardW, cardH = 0.37, 1.2, 7.8, 5.6

    # Card background
    add_shape(slide, MSO_SHAPE.RECTANGLE, cardX, cardY, cardW, cardH,
              fill_color=C["lightGrey"], line_color=C["medGrey"], line_width=0.5)

    # Description section header
    add_shape(slide, MSO_SHAPE.RECTANGLE, cardX, cardY, cardW, 0.38,
              fill_color=C["infTeal"])
    add_textbox(slide, cardX + 0.15, cardY + 0.03, cardW - 0.3, 0.32,
                "Description", font_size=11, bold=True, color=C["white"])

    # Description body
    if isinstance(description, list):
        add_rich_text(slide, cardX + 0.15, cardY + 0.44, cardW - 0.3, 1.5,
                      description, font_size=10)
    else:
        add_textbox(slide, cardX + 0.15, cardY + 0.44, cardW - 0.3, 1.5,
                    description, font_size=10, color=C["black"])

    # Status section header
    statusY = cardY + 2.0
    add_shape(slide, MSO_SHAPE.RECTANGLE, cardX, statusY, cardW, 0.35,
              fill_color=C["infTeal"])
    add_textbox(slide, cardX + 0.15, statusY + 0.02, cardW - 0.3, 0.31,
                "Status", font_size=11, bold=True, color=C["white"])

    # Status badge
    add_shape(slide, MSO_SHAPE.RECTANGLE, cardX + 0.15, statusY + 0.42, 2.2, 0.32,
              fill_color=status_color)
    add_textbox(slide, cardX + 0.15, statusY + 0.42, 2.2, 0.32,
                status, font_size=9, bold=True, color=C["white"],
                alignment=PP_ALIGN.CENTER)

    # Status details
    if status_details:
        if isinstance(status_details, list):
            add_rich_text(slide, cardX + 0.15, statusY + 0.82, cardW - 0.3, 1.1,
                          status_details, font_size=9)
        else:
            add_textbox(slide, cardX + 0.15, statusY + 0.82, cardW - 0.3, 1.1,
                        status_details, font_size=9, color=C["black"])

    # Users section
    if users:
        usersY = cardY + 4.0
        add_shape(slide, MSO_SHAPE.RECTANGLE, cardX, usersY, cardW, 0.35,
                  fill_color=C["infTeal"])
        add_textbox(slide, cardX + 0.15, usersY + 0.02, cardW - 0.3, 0.31,
                    "Actual Users", font_size=11, bold=True, color=C["white"])

        user_items = [{"text": f"• {u}", "color": C["black"]} for u in users]
        add_rich_text(slide, cardX + 0.15, usersY + 0.4, cardW - 0.3, 1.0,
                      user_items, font_size=9)

    # ─── Right Panel: Team + Metrics ───
    rightX = 8.5

    # Team placeholder circles
    add_textbox(slide, rightX, 1.3, 1.5, 0.3, "Team",
                font_size=11, bold=True, color=C["infTeal"])

    photo_colors = ["D6EAF8", "C8E6C9", "FFE0B2", "E1BEE7", "B3E5FC"]
    for i in range(5):
        px = rightX + (i % 4) * 0.65
        py = 1.7 + (i // 4) * 0.65
        add_shape(slide, MSO_SHAPE.OVAL, px, py, 0.55, 0.55,
                  fill_color=photo_colors[i], line_color=C["medGrey"], line_width=0.5)

    # Metrics box
    if metrics:
        savingsY = 4.6
        savingsW = 4.5
        add_shape(slide, MSO_SHAPE.RECTANGLE, rightX - 0.1, savingsY, savingsW, 2.0,
                  fill_color=C["lightGrey"], line_color=C["medGrey"], line_width=0.5)
        add_shape(slide, MSO_SHAPE.RECTANGLE, rightX - 0.1, savingsY, savingsW, 0.35,
                  fill_color=C["infTeal"])
        add_textbox(slide, rightX - 0.1, savingsY, savingsW, 0.35,
                    "Expected Savings", font_size=10, bold=True, color=C["white"],
                    alignment=PP_ALIGN.CENTER)

        metric_count = len(metrics)
        metricW = (savingsW - 0.4) / metric_count
        for i, metric in enumerate(metrics):
            metricX = rightX + 0.1 + i * metricW
            add_textbox(slide, metricX, savingsY + 0.4, metricW, 0.28,
                        metric["label"], font_size=9, color=C["grey"],
                        alignment=PP_ALIGN.CENTER)
            add_textbox(slide, metricX, savingsY + 0.65, metricW, 1.1,
                        metric["value"], font_size=42, bold=True,
                        color=metric["color"], alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# COMPARISON TABLE LAYOUT
# ═══════════════════════════════════════════════════════════

def create_comparison_table(slide, *,
                            title,
                            headers,
                            rows,
                            col_widths=None,
                            logo_path=None,
                            page_number=None):
    """
    Create a comparison table slide.

    Args:
        slide: slide object
        title (str): slide title
        headers (list[str]): column headers
        rows (list[list[str]]): table data, each inner list is a row
        col_widths (list[float], optional): column widths in inches
        logo_path (str, optional): path to logo PNG
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]
    ct = L["content"]

    add_infineon_header(slide, title, logo_path=logo_path)
    add_infineon_footer(slide, page_number=page_number)

    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header row

    if col_widths is None:
        col_widths = [ct["w"] / num_cols] * num_cols

    # Create table
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(ct["x"]), Inches(ct["y"]),
        Inches(ct["w"]), Inches(min(num_rows * 0.5, ct["h"]))
    )
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Header row
    for j, header_text in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header_text
        _style_cell(cell, fill_color=C["infTeal"], font_color=C["white"],
                    font_size=10, bold=True, alignment=PP_ALIGN.CENTER)

    # Data rows
    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = cell_text
            if j == 0:
                _style_cell(cell, fill_color="E0F2F1",
                            font_color=C["infTeal"], font_size=9, bold=True)
            else:
                bg = C["white"] if i % 2 == 0 else C["lightGrey"]
                _style_cell(cell, fill_color=bg,
                            font_color=C["black"], font_size=9)


def _style_cell(cell, fill_color, font_color, font_size=9,
                bold=False, alignment=PP_ALIGN.LEFT):
    """Style a table cell."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = hex_to_rgb(fill_color)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = hex_to_rgb(font_color)
        paragraph.font.bold = bold
        paragraph.font.name = "Arial"
        paragraph.alignment = alignment


# ═══════════════════════════════════════════════════════════
# ARCHITECTURE DIAGRAM LAYOUT
# ═══════════════════════════════════════════════════════════

def create_architecture_diagram(slide, *,
                                title,
                                boxes,
                                subtitle=None,
                                callout_text=None,
                                logo_path=None,
                                page_number=None):
    """
    Create an architecture flow diagram with boxes and arrows.

    Args:
        slide: slide object
        title (str): slide title
        boxes (list[dict]): each dict has keys: label, color, subtext (optional)
        subtitle (str, optional): subtitle text
        callout_text (str, optional): explanatory text in a callout box
        logo_path (str, optional): path to logo PNG
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]

    add_infineon_header(slide, title, subtitle, logo_path)
    add_infineon_footer(slide, page_number=page_number)

    boxH, boxW, boxY = 1.3, 2.1, 2.0
    spacing = 0.4
    totalW = len(boxes) * boxW + (len(boxes) - 1) * spacing
    startX = (L["slideWidth"] - totalW) / 2

    for i, box in enumerate(boxes):
        bx = startX + i * (boxW + spacing)

        # Box
        add_shape(slide, MSO_SHAPE.RECTANGLE, bx, boxY, boxW, boxH,
                  fill_color=box["color"])

        # Label
        add_textbox(slide, bx, boxY, boxW, boxH, box["label"],
                    font_size=12, bold=True, color=C["white"],
                    alignment=PP_ALIGN.CENTER)

        # Subtext
        if box.get("subtext"):
            add_textbox(slide, bx, boxY + boxH + 0.12, boxW, 0.8,
                        box["subtext"], font_size=8, color=C["grey"],
                        alignment=PP_ALIGN.CENTER)

        # Arrow to next box
        if i < len(boxes) - 1:
            arrowX = bx + boxW
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      arrowX, boxY + boxH / 2 - 0.04, spacing, 0.08,
                      fill_color=C["grey"])
            add_textbox(slide, arrowX + spacing - 0.15,
                        boxY + boxH / 2 - 0.16, 0.24, 0.32,
                        "▶", font_size=12, color=C["grey"],
                        alignment=PP_ALIGN.LEFT)

    # Callout box
    if callout_text:
        callY = 5.0
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.37, callY, 12.6, 1.6,
                  fill_color="E0F2F1", line_color=C["infTeal"], line_width=0.75)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.37, callY, 0.7, 1.6,
                  fill_color=C["infTeal"])
        add_textbox(slide, 0.37, callY, 0.7, 1.6, "ℹ",
                    font_size=22, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, 1.2, callY + 0.1, 11.5, 1.4,
                    callout_text, font_size=10, color=C["black"])


# ═══════════════════════════════════════════════════════════
# ROADMAP TIMELINE LAYOUT
# ═══════════════════════════════════════════════════════════

def create_roadmap(slide, *,
                   title,
                   time_labels,
                   tracks,
                   subtitle=None,
                   logo_path=None,
                   page_number=None):
    """
    Create a roadmap / timeline slide.

    Args:
        slide: slide object
        title (str): slide title
        time_labels (list[str]): e.g. ["Q1 2025", "Q2 2025", ...]
        tracks (list[dict]): each dict has:
            - name (str): track name
            - phases (list[dict]): each with label, color, start, duration
        subtitle (str, optional)
        logo_path (str, optional)
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]

    add_infineon_header(slide, title, subtitle, logo_path)
    add_infineon_footer(slide, page_number=page_number)

    timeStartX = 1.5
    timeWidth = (L["slideWidth"] - timeStartX - 0.5) / len(time_labels)
    trackStartY = 1.8
    trackHeight = 0.9
    trackSpacing = 0.3

    # Time headers
    for i, label in enumerate(time_labels):
        x = timeStartX + i * timeWidth
        add_textbox(slide, x, 1.35, timeWidth - 0.1, 0.35, label,
                    font_size=10, bold=True, color=C["infTeal"],
                    alignment=PP_ALIGN.CENTER)
        # Grid line
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  x + timeWidth - 0.05, 1.3, 0.01, 5.3,
                  fill_color="E0E0E0")

    # Tracks
    for ti, track in enumerate(tracks):
        rowY = trackStartY + ti * (trackHeight + trackSpacing)

        # Track name
        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.37, rowY, 1.0, trackHeight,
                  fill_color=C["infTeal"])
        add_textbox(slide, 0.37, rowY, 1.0, trackHeight, track["name"],
                    font_size=8, bold=True, color=C["white"],
                    alignment=PP_ALIGN.CENTER)

        # Track background
        total_track_w = len(time_labels) * timeWidth
        fill = C["lightGrey"] if ti % 2 == 0 else C["white"]
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  timeStartX, rowY, total_track_w, trackHeight,
                  fill_color=fill, line_color="ECECEC", line_width=0.3)

        # Phase bars
        for phase in track["phases"]:
            phaseX = timeStartX + phase["start"] * timeWidth
            phaseW = phase["duration"] * timeWidth
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      phaseX, rowY + 0.15, phaseW, 0.6,
                      fill_color=phase["color"])
            add_textbox(slide, phaseX, rowY + 0.15, phaseW, 0.6,
                        phase["label"], font_size=8, bold=True,
                        color=C["white"], alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# TWO-COLUMN LAYOUT
# ═══════════════════════════════════════════════════════════

def create_two_column_layout(slide, *,
                             title,
                             left_panel,
                             right_panel,
                             logo_path=None,
                             page_number=None):
    """
    Create a two-column layout slide.

    Args:
        slide: slide object
        title (str): slide title
        left_panel (dict): keys: title, bg_color, items (list of {title, text})
        right_panel (dict): keys: title, bg_color, items (list of {icon, label})
        logo_path (str, optional)
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]

    add_infineon_header(slide, title, logo_path=logo_path)
    add_infineon_footer(slide, page_number=page_number)

    panelW, panelH, panelY = 6.1, 5.5, 1.3

    # ── Left panel ──
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.37, panelY, panelW, panelH,
              fill_color=left_panel.get("bg_color", C["infTeal"]))

    add_textbox(slide, 0.55, panelY + 0.2, panelW - 0.36, 0.8,
                left_panel["title"], font_size=18, bold=True,
                color=C["white"], alignment=PP_ALIGN.CENTER)

    for i, item in enumerate(left_panel["items"]):
        itemY = panelY + 1.2 + i * 2.1

        add_shape(slide, MSO_SHAPE.RECTANGLE, 0.55, itemY, panelW - 0.36, 0.35,
                  fill_color=C["infOrange"])
        add_textbox(slide, 0.55, itemY, panelW - 0.36, 0.35,
                    item["title"], font_size=10, bold=True,
                    color=C["white"], alignment=PP_ALIGN.CENTER)

        if item.get("text"):
            add_textbox(slide, 0.65, itemY + 0.4, panelW - 0.56, 1.55,
                        item["text"], font_size=9, color=C["white"])

    # ── Right panel ──
    rightX = 0.37 + panelW + 0.25
    add_shape(slide, MSO_SHAPE.RECTANGLE, rightX, panelY, panelW, panelH,
              fill_color=right_panel.get("bg_color", "E0F2F1"),
              line_color=C["infTeal"], line_width=0.75)

    add_textbox(slide, rightX + 0.2, panelY + 0.2, panelW - 0.4, 0.55,
                right_panel["title"], font_size=17, bold=True,
                color=C["infTeal"], alignment=PP_ALIGN.CENTER)

    for i, item in enumerate(right_panel["items"]):
        itemY = panelY + 1.0 + i * 0.72
        fill = C["white"] if i % 2 == 0 else C["lightGrey"]
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  rightX + 0.2, itemY, panelW - 0.4, 0.6,
                  fill_color=fill, line_color="E0E0E0", line_width=0.3)

        if item.get("icon"):
            add_textbox(slide, rightX + 0.25, itemY, 0.5, 0.6,
                        item["icon"], font_size=14,
                        alignment=PP_ALIGN.CENTER)

        add_textbox(slide, rightX + 0.8, itemY, panelW - 1.0, 0.6,
                    item["label"], font_size=11, color=C["black"])


# ═══════════════════════════════════════════════════════════
# METRIC / KPI DASHBOARD LAYOUT
# ═══════════════════════════════════════════════════════════

def create_metric_dashboard(slide, *,
                            title,
                            metrics,
                            logo_path=None,
                            page_number=None):
    """
    Create a KPI / metrics dashboard slide with large numbers.

    Args:
        slide: slide object
        title (str): slide title
        metrics (list[dict]): each dict has:
            - label (str): description
            - value (str): the big number / percentage
            - color (str): hex color for the value
            - subtitle (str, optional): extra context below value
        logo_path (str, optional)
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]

    add_infineon_header(slide, title, logo_path=logo_path)
    add_infineon_footer(slide, page_number=page_number)

    count = len(metrics)
    card_spacing = 0.3
    total_w = L["content"]["w"]
    card_w = (total_w - (count - 1) * card_spacing) / count
    card_h = 4.0
    startX = L["content"]["x"]
    startY = L["content"]["y"] + 0.3

    for i, m in enumerate(metrics):
        cx = startX + i * (card_w + card_spacing)

        # Card background
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  cx, startY, card_w, card_h,
                  fill_color=C["lightGrey"], line_color=C["medGrey"], line_width=0.5)

        # Color accent bar at top
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  cx, startY, card_w, 0.1,
                  fill_color=m["color"])

        # Label
        add_textbox(slide, cx + 0.2, startY + 0.3, card_w - 0.4, 0.5,
                    m["label"], font_size=12, bold=True, color=C["grey"],
                    alignment=PP_ALIGN.CENTER)

        # Value
        add_textbox(slide, cx + 0.2, startY + 1.0, card_w - 0.4, 2.0,
                    m["value"], font_size=48, bold=True,
                    color=m["color"], alignment=PP_ALIGN.CENTER)

        # Subtitle
        if m.get("subtitle"):
            add_textbox(slide, cx + 0.2, startY + 3.2, card_w - 0.4, 0.6,
                        m["subtitle"], font_size=9, color=C["grey"],
                        alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# IMAGE GALLERY LAYOUT
# ═══════════════════════════════════════════════════════════

def create_image_gallery(slide, *,
                         title,
                         images,
                         logo_path=None,
                         page_number=None):
    """
    Create an image gallery slide.

    Args:
        slide: slide object
        title (str): slide title
        images (list[dict]): each dict has:
            - path (str): path to the image file
            - caption (str, optional): caption text below image
        logo_path (str, optional)
    """
    import os
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]

    add_infineon_header(slide, title, logo_path=logo_path)
    add_infineon_footer(slide, page_number=page_number)

    count = len(images)
    spacing = 0.3
    total_w = L["content"]["w"]
    img_w = (total_w - (count - 1) * spacing) / count
    img_h = 3.5
    startX = L["content"]["x"]
    startY = L["content"]["y"] + 0.3

    for i, img in enumerate(images):
        ix = startX + i * (img_w + spacing)

        if os.path.isfile(img["path"]):
            slide.shapes.add_picture(
                img["path"],
                Inches(ix), Inches(startY),
                Inches(img_w), Inches(img_h)
            )
        else:
            # Placeholder
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      ix, startY, img_w, img_h,
                      fill_color=C["lightGrey"], line_color=C["medGrey"],
                      line_width=0.5)
            add_textbox(slide, ix, startY + img_h / 2 - 0.2,
                        img_w, 0.4, f"[{img['path']}]",
                        font_size=8, color=C["grey"],
                        alignment=PP_ALIGN.CENTER)

        if img.get("caption"):
            add_textbox(slide, ix, startY + img_h + 0.1, img_w, 0.4,
                        img["caption"], font_size=9, color=C["grey"],
                        alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# KPI CHART (Planned vs Achieved)
# ═══════════════════════════════════════════════════════════

def create_kpi_chart(slide, *,
                     title,
                     categories,
                     planned,
                     achieved,
                     chart_title="Issues - Planned vs Achieved",
                     planned_color=None,
                     achieved_color=None,
                     logo_path=None,
                     page_number=None):
    """
    Create a KPI chart slide with clustered columns (Planned vs Achieved).

    Args:
        slide: slide object
        title (str): slide title
        categories (list[str]): category labels (e.g. quarters, months, sprints)
        planned (list[int|float]): planned values per category
        achieved (list[int|float]): achieved values per category
        chart_title (str, optional): chart heading text
        planned_color (str, optional): hex color for planned series (default: infTeal)
        achieved_color (str, optional): hex color for achieved series (default: infOrange)
        logo_path (str, optional): path to logo PNG
        page_number (int, optional): page number for footer
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]

    add_infineon_header(slide, title, logo_path=logo_path)
    add_infineon_footer(slide, page_number=page_number)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Planned", planned)
    chart_data.add_series("Achieved", achieved)

    ct = L["content"]
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(ct["x"]), Inches(ct["y"]),
        Inches(ct["w"]), Inches(ct["h"]),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    plot = chart.plots[0]
    plot.gap_width = 100

    # Planned series
    series_planned = plot.series[0]
    series_planned.format.fill.solid()
    series_planned.format.fill.fore_color.rgb = hex_to_rgb(planned_color or C["infTeal"])

    # Achieved series
    series_achieved = plot.series[1]
    series_achieved.format.fill.solid()
    series_achieved.format.fill.fore_color.rgb = hex_to_rgb(achieved_color or C["infOrange"])

    # Data labels
    for series in plot.series:
        series.has_data_labels = True
        series.data_labels.font.size = Pt(8)
        series.data_labels.number_format = "0"


# ═══════════════════════════════════════════════════════════
# EFFORT CHART (Planned vs Delivered in days)
# ═══════════════════════════════════════════════════════════

def create_effort_chart(slide, *,
                        title,
                        categories,
                        planned_days,
                        achieved_days,
                        chart_title="Effort (days) - Planned vs Achieved",
                        planned_color=None,
                        achieved_color=None,
                        logo_path=None,
                        page_number=None):
    """
    Create an effort trend chart slide with clustered columns.

    Args:
        slide: slide object
        title (str): slide title
        categories (list[str]): category labels (e.g. quarters, months, sprints)
        planned_days (list[float]): planned effort in days per category
        achieved_days (list[float]): achieved effort in days per category
        chart_title (str, optional): chart heading text
        planned_color (str, optional): hex color for planned series (default: blue)
        achieved_color (str, optional): hex color for achieved series (default: lime)
        logo_path (str, optional): path to logo PNG
        page_number (int, optional): page number for footer
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]

    add_infineon_header(slide, title, logo_path=logo_path)
    add_infineon_footer(slide, page_number=page_number)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Planned (days)", planned_days)
    chart_data.add_series("Achieved (days)", achieved_days)

    ct = L["content"]
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(ct["x"]), Inches(ct["y"]),
        Inches(ct["w"]), Inches(ct["h"]),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    plot = chart.plots[0]
    plot.gap_width = 100

    series_planned = plot.series[0]
    series_planned.format.fill.solid()
    series_planned.format.fill.fore_color.rgb = hex_to_rgb(planned_color or C["blue"])

    series_achieved = plot.series[1]
    series_achieved.format.fill.solid()
    series_achieved.format.fill.fore_color.rgb = hex_to_rgb(achieved_color or C["lime"])

    for series in plot.series:
        series.has_data_labels = True
        series.data_labels.font.size = Pt(8)
        series.data_labels.number_format = "0.0"


# ═══════════════════════════════════════════════════════════
# BURNDOWN CHART (Stacked columns)
# ═══════════════════════════════════════════════════════════

def create_burndown_chart(slide, *,
                          title,
                          categories,
                          completed,
                          remaining,
                          added,
                          show_velocity=True,
                          velocity_unit="items/period",
                          logo_path=None,
                          page_number=None):
    """
    Create a burndown stacked chart slide.

    Args:
        slide: slide object
        title (str): slide title
        categories (list[str]): category labels (e.g. sprint names, months, weeks)
        completed (list[int|float]): completed work per category
        remaining (list[int|float]): remaining work per category
        added (list[int|float]): added work per category
        show_velocity (bool): if True, show average velocity annotation
        velocity_unit (str): unit label for velocity (e.g. "items/sprint", "tasks/week")
        logo_path (str, optional): path to logo PNG
        page_number (int, optional): page number for footer
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]

    add_infineon_header(slide, title, logo_path=logo_path)
    add_infineon_footer(slide, page_number=page_number)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Completed", completed)
    chart_data.add_series("Remaining", remaining)
    chart_data.add_series("Added", added)

    ct = L["content"]
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(ct["x"]), Inches(ct["y"]),
        Inches(ct["w"]), Inches(ct["h"]),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    plot = chart.plots[0]

    # Completed - green
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = hex_to_rgb(C["lime"])

    # Remaining - dark teal
    plot.series[1].format.fill.solid()
    plot.series[1].format.fill.fore_color.rgb = hex_to_rgb(C["infTeal"])

    # Added - blue
    plot.series[2].format.fill.solid()
    plot.series[2].format.fill.fore_color.rgb = hex_to_rgb(C["blue"])

    # Velocity annotation
    if show_velocity and len(completed) > 1:
        avg_velocity = sum(completed) / len(completed)
        add_textbox(slide, ct["x"], ct["y"] + ct["h"] + 0.05,
                    5, 0.3,
                    f"Average Velocity: {avg_velocity:.1f} {velocity_unit}",
                    font_size=9, color=C["grey"])


# ═══════════════════════════════════════════════════════════
# CUMULATIVE TREND CHART (Line chart)
# ═══════════════════════════════════════════════════════════

def create_cumulative_chart(slide, *,
                            title,
                            categories,
                            cumulative_planned,
                            cumulative_achieved,
                            planned_color=None,
                            achieved_color=None,
                            logo_path=None,
                            page_number=None):
    """
    Create a cumulative trend line chart slide.

    Args:
        slide: slide object
        title (str): slide title
        categories (list[str]): category labels (e.g. quarters, months, sprints)
        cumulative_planned (list[int|float]): cumulative planned counts
        cumulative_achieved (list[int|float]): cumulative achieved counts
        planned_color (str, optional): hex color for planned line (default: infTeal)
        achieved_color (str, optional): hex color for achieved line (default: infOrange)
        logo_path (str, optional): path to logo PNG
        page_number (int, optional): page number for footer
    """
    C = INFINEON_THEME["colors"]
    L = INFINEON_THEME["layout"]

    add_infineon_header(slide, title, logo_path=logo_path)
    add_infineon_footer(slide, page_number=page_number)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Planned (cumulative)", cumulative_planned)
    chart_data.add_series("Achieved (cumulative)", cumulative_achieved)

    ct = L["content"]
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(ct["x"]), Inches(ct["y"]),
        Inches(ct["w"]), Inches(ct["h"]),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    # Style lines
    series_planned = chart.series[0]
    series_planned.format.line.color.rgb = hex_to_rgb(planned_color or C["infTeal"])
    series_planned.format.line.width = Pt(2.5)
    series_planned.smooth = False

    series_achieved = chart.series[1]
    series_achieved.format.line.color.rgb = hex_to_rgb(achieved_color or C["infOrange"])
    series_achieved.format.line.width = Pt(2.5)
    series_achieved.smooth = False

    for series in [series_planned, series_achieved]:
        series.has_data_labels = True
        series.data_labels.font.size = Pt(8)
        series.data_labels.number_format = "0"
        series.data_labels.position = XL_LABEL_POSITION.ABOVE
