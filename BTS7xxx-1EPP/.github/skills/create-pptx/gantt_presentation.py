#!/usr/bin/env python3
"""
Agile Scrum Rollout — Gantt Chart PowerPoint
Infineon Technologies Brand Style (2023-2)

Reads Agile_Rollout_Gantt.xlsx and generates a multi-slide PPTX with
per-phase Gantt charts in Infineon visual style.

Run from the ifx_ppt folder:
    python gantt_presentation.py
"""

import sys
import os

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from infineon_theme import (
    INFINEON_THEME, hex_to_rgb,
    create_presentation, create_title_slide,
    find_infineon_logo, add_textbox, add_shape,
    add_infineon_header, add_infineon_footer,
)
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

C = INFINEON_THEME["colors"]

# ──────────────────────────────────────────────────────────────────
# TIMELINE COLUMN MAP  (0-based: 0 = Apr 2026, 25 = May 2028)
# ──────────────────────────────────────────────────────────────────
MONTHS = [
    ("Apr", "2026"), ("May", "2026"), ("Jun", "2026"), ("Jul", "2026"),
    ("Aug", "2026"), ("Sep", "2026"), ("Oct", "2026"), ("Nov", "2026"), ("Dec", "2026"),
    ("Jan", "2027"), ("Feb", "2027"), ("Mar", "2027"), ("Apr", "2027"), ("May", "2027"),
    ("Jun", "2027"), ("Jul", "2027"), ("Aug", "2027"), ("Sep", "2027"), ("Oct", "2027"),
    ("Nov", "2027"), ("Dec", "2027"),
    ("Jan", "2028"), ("Feb", "2028"), ("Mar", "2028"), ("Apr", "2028"), ("May", "2028"),
]

# ──────────────────────────────────────────────────────────────────
# PHASE COLOR PALETTE  (Infineon 2023-2 official accent palette)
#   accent1 infTeal   0A8276  → Phase 0
#   accent3 infOrange F97414  → Phase 1
#   accent4 infLime   9BBA43  → Phase 2
#   accent5 infYellow FCD442  → Phase 3 Wave 2
#           blue      0070C0  → Phase 3 Wave 3
#   accent6 infMagenta 9C216E → Phase 4
#   dk1     textPrimary 1D1D1D → milestones & header labels
# ──────────────────────────────────────────────────────────────────
PC = {
    "P0":   {"bar": "0A8276", "light": "6CB4AD", "hdr": "1D1D1D"},   # infTeal   (accent1)
    "P1":   {"bar": "F97414", "light": "FF9737", "hdr": "1D1D1D"},   # infOrange (accent3)
    "P2":   {"bar": "9BBA43", "light": "6CB4AD", "hdr": "1D1D1D"},   # infLime   (accent4)
    "P3W2": {"bar": "FCD442", "light": "FBE29B", "hdr": "1D1D1D"},   # infYellow (accent5)
    "P3W3": {"bar": "0070C0", "light": "E3A9CB", "hdr": "1D1D1D"},   # Infineon blue
    "P4":   {"bar": "9C216E", "light": "E3A9CB", "hdr": "1D1D1D"},   # infMagenta (accent6)
}

MILESTONE_COLOR = "1D1D1D"   # Infineon textPrimary (dk1)

# ──────────────────────────────────────────────────────────────────
# ACTIVITY DEFINITIONS
# Tuple: (phase_key, label, bar_start, bar_end, is_milestone, is_phase_hdr, bar_color_override)
#   bar_start/bar_end: 0-based month index (0=Apr2026 … 25=May2028)
#   bar_color_override: hex string for lighter "sub-activity" bars, or None
# ──────────────────────────────────────────────────────────────────
ACTIVITIES = [
    # Phase 0 · Foundation ─────────────────────────────────────────
    ("P0",   "Phase 0 · Foundation",                   0,  1, False, True,  None),
    ("P0",   "Executive Alignment & Budget Approval",   0,  0, False, False, None),
    ("P0",   "Form Agile Transformation Office (ATO)",  0,  0, False, False, None),
    ("P0",   "Consulting Partner RFP & Selection",      0,  1, False, False, None),
    ("P0",   "Current State Assessment",                0,  1, False, False, None),
    ("P0",   "◆ Transformation Roadmap Presented",      1,  1, True,  False, None),

    # Phase 1 · Design ─────────────────────────────────────────────
    ("P1",   "Phase 1 · Design",                        1,  3, False, True,  None),
    ("P1",   "Scaling Framework Design (SAFe / ARTs)",  1,  2, False, False, None),
    ("P1",   "Work Item Mapping (Epics / Features)",    1,  2, False, False, None),
    ("P1",   "Definition of Done per Team Type",        2,  2, False, False, None),
    ("P1",   "Waterfall Interface Design",              2,  2, False, False, None),
    ("P1",   "JIRA / Confluence Tooling Setup",         2,  3, False, False, None),
    ("P1",   "ISO 26262 / Safety Alignment",            1,  2, False, False, None),
    ("P1",   "◆ Architecture & Design Complete",        3,  3, True,  False, None),

    # Phase 2 · Pilot Wave 1 ────────────────────────────────────────
    ("P2",   "Phase 2 · Pilot (Wave 1)",                2,  8, False, True,  None),
    ("P2",   "Pilot Team Selection & Preparation",      2,  3, False, False, None),
    ("P2",   "Story Mapping Workshops",                 3,  3, False, False, None),
    ("P2",   "Pilot Team Backlog Creation",             3,  3, False, False, None),
    ("P2",   "Sprints 1–2  (Intensively Coached)",      3,  4, False, False, None),
    ("P2",   "Sprints 3–4  (Coach Tapering)",           4,  5, False, False, None),
    ("P2",   "Sprints 5–6",                             5,  6, False, False, None),
    ("P2",   "Sprints 7–8  (Internal Coaches Leading)", 6,  7, False, False, None),
    ("P2",   "Sprints 9–10",                            7,  8, False, False, None),
    ("P2",   "Internal Coach Training (SM / ICP-ACC)",  4,  6, False, False, "6CB4AD"),  # tealLight
    ("P2",   "◆ First PI Planning Event (Pilot ARTs)",  6,  6, True,  False, None),
    ("P2",   "◆ Pilot Retrospective & Go/No-Go",        8,  8, True,  False, None),

    # Phase 3 · Scale — Wave 2 ─────────────────────────────────────
    ("P3W2", "Phase 3 · Scale — Wave 2",                8, 15, False, True,  None),
    ("P3W2", "Wave 2 ART Onboarding (4 ARTs)",          9, 11, False, False, None),
    ("P3W2", "Wave 2 Story Mapping & Backlog",          9, 10, False, False, None),
    ("P3W2", "Wave 2 First Sprints (Coached)",         11, 12, False, False, None),
    ("P3W2", "Wave 2 Sprints Continue",                12, 14, False, False, None),
    ("P3W2", "Internal RTEs Designated & Trained",      9, 11, False, False, "FBE29B"),  # yellowSoft
    ("P3W2", "Communities of Practice Established",    11, 11, False, False, "FBE29B"),  # yellowSoft
    ("P3W2", "◆ PI Planning — Wave 1+2 Combined",      11, 11, True,  False, None),
    ("P3W2", "◆ Consulting Firm → Advisory Only",      15, 15, True,  False, None),

    # Phase 3 · Scale — Wave 3 ─────────────────────────────────────
    ("P3W3", "Phase 3 · Scale — Wave 3",               15, 21, False, True,  None),
    ("P3W3", "Wave 3 ART Onboarding (~55 teams)",      15, 17, False, False, None),
    ("P3W3", "Wave 3 Backlog & First Sprints",         17, 19, False, False, None),
    ("P3W3", "Wave 3 Full Sprint Cadence",             19, 22, False, False, None),
    ("P3W3", "Inspect & Adapt Events (All ARTs)",      18, 22, False, False, "E3A9CB"),  # pinkLight
    ("P3W3", "◆ All ARTs in PI Planning Cadence",      21, 21, True,  False, None),

    # Phase 4 · Sustain & Mature ───────────────────────────────────
    ("P4",   "Phase 4 · Sustain & Mature",             15, 24, False, True,  None),
    ("P4",   "CoPs Self-Sustaining",                   17, 20, False, False, None),
    ("P4",   "Metrics & Portfolio Dashboard Live",     15, 16, False, False, None),
    ("P4",   "◆ Maturity Assessment #1",               18, 18, True,  False, None),
    ("P4",   "◆ Consulting Firm Exit",                 22, 22, True,  False, None),
    ("P4",   "Transformation Complete",                22, 23, False, False, None),
    ("P4",   "◆ Maturity Assessment #2",               24, 24, True,  False, None),
]


# ──────────────────────────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────────────────────────

def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = hex_to_rgb(C["white"])
    return slide


def _fs(row_h):
    """Pick font size based on row height (inches)."""
    if row_h >= 0.38:
        return 8
    if row_h >= 0.28:
        return 7
    return 6


# ──────────────────────────────────────────────────────────────────
# CORE GANTT DRAWING  (one slide segment)
# ──────────────────────────────────────────────────────────────────

def draw_gantt_slide(slide, title, activities, col_start, col_end,
                     page_number=None, logo_path=None, reserve_bottom=0.0):
    """
    Paint a Gantt chart segment on an existing blank slide.

    Parameters
    ----------
    slide        : pptx slide object
    title        : string — displayed as slide header
    activities   : list of ACTIVITIES tuples to render
    col_start    : first month column (0-based, inclusive)
    col_end      : last  month column (0-based, inclusive)
    page_number  : int or None
    logo_path    : path to logo PNG or None
    reserve_bottom : extra inches to leave free below the chart
                     (used when legend will be drawn below)
    """
    add_infineon_header(slide, title, logo_path=logo_path)
    add_infineon_footer(slide, page_number=page_number)

    N      = col_end - col_start + 1
    n_rows = len(activities)

    # ── Geometry ─────────────────────────────────────────────────
    MARGIN   = 0.35
    LABEL_W  = 2.60
    TOP_Y    = 1.08
    BOT_Y    = 7.04 - reserve_bottom
    CHART_X  = MARGIN + LABEL_W
    CHART_W  = 13.333 - MARGIN * 2 - LABEL_W
    TOTAL_H  = BOT_Y - TOP_Y

    YEAR_H   = 0.24
    MONTH_H  = 0.20
    DATA_H   = TOTAL_H - YEAR_H - MONTH_H

    ROW_H    = min(DATA_H / n_rows, 0.52)
    COL_W    = CHART_W / N

    YEAR_Y   = TOP_Y
    MONTH_Y  = YEAR_Y + YEAR_H
    DATA_Y   = MONTH_Y + MONTH_H

    FS = _fs(ROW_H)
    BAR_INSET = max(0.04, ROW_H * 0.14)

    # ── Label-column top-left header cell ────────────────────────
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              MARGIN, YEAR_Y, LABEL_W, YEAR_H + MONTH_H,
              fill_color=C["infTeal"], line_color=None)
    add_textbox(slide, MARGIN + 0.08, YEAR_Y, LABEL_W - 0.12, YEAR_H + MONTH_H,
                "Phase / Activity",
                font_size=8, bold=True, color=C["white"],
                alignment=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)

    # ── Year header bands ─────────────────────────────────────────
    # Colors: infTeal for 2026, textPrimary for 2027, darkGrey for 2028
    YEAR_FILL = {"2026": C["infTeal"], "2027": C["textPrimary"], "2028": C["darkGrey"]}
    cur_yr = None
    yr_start_i = 0

    def _flush_year(i_s, i_e, yr):
        x = CHART_X + i_s * COL_W
        w = (i_e - i_s + 1) * COL_W
        yc = YEAR_FILL.get(yr, C["infTeal"])
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, YEAR_Y, w, YEAR_H,
                  fill_color=yc, line_color="FFFFFF", line_width=0.3)
        add_textbox(slide, x + 0.02, YEAR_Y, w - 0.04, YEAR_H,
                    yr, font_size=8, bold=True, color=C["white"],
                    alignment=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    for i in range(N):
        _, yr = MONTHS[col_start + i]
        if yr != cur_yr:
            if cur_yr is not None:
                _flush_year(yr_start_i, i - 1, cur_yr)
            cur_yr = yr
            yr_start_i = i
    _flush_year(yr_start_i, N - 1, cur_yr)

    # ── Month header cells ────────────────────────────────────────
    for i in range(N):
        mo, _ = MONTHS[col_start + i]
        x = CHART_X + i * COL_W
        fill = C["lightGrey"] if i % 2 == 0 else C["medGrey"]
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, MONTH_Y, COL_W, MONTH_H,
                  fill_color=fill, line_color="FFFFFF", line_width=0.25)
        mo_fs = max(5, min(8, int(COL_W * 18)))
        add_textbox(slide, x + 0.01, MONTH_Y, COL_W - 0.02, MONTH_H,
                    mo, font_size=mo_fs, color=C["textPrimary"],
                    alignment=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    # ── Row backgrounds ───────────────────────────────────────────
    total_row_w = LABEL_W + CHART_W
    for r, act in enumerate(activities):
        phase_key, label, bar_start, bar_end, is_ms, is_hdr, _ = act
        y = DATA_Y + r * ROW_H
        if is_hdr:
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      MARGIN, y, total_row_w, ROW_H,
                      fill_color=PC[phase_key]["hdr"], line_color="FFFFFF", line_width=0.2)
        else:
            row_bg = C["lightGrey"] if r % 2 == 0 else C["white"]
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      MARGIN, y, total_row_w, ROW_H,
                      fill_color=row_bg, line_color=C["medGrey"], line_width=0.15)

    # ── Vertical grid lines ───────────────────────────────────────
    grid_h = ROW_H * n_rows
    for i in range(N + 1):
        x = CHART_X + i * COL_W
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, DATA_Y, 0.003, grid_h,
                  fill_color=C["medGrey"], line_color=None)

    # ── Activity labels ───────────────────────────────────────────
    for r, act in enumerate(activities):
        phase_key, label, bar_start, bar_end, is_ms, is_hdr, _ = act
        y = DATA_Y + r * ROW_H
        if is_hdr:
            add_textbox(slide, MARGIN + 0.08, y, LABEL_W - 0.12, ROW_H,
                        label, font_size=min(FS + 1, 9), bold=True, color=C["white"],
                        alignment=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
        else:
            lbl_color = MILESTONE_COLOR if is_ms else "1D1D1D"
            add_textbox(slide, MARGIN + 0.10, y, LABEL_W - 0.13, ROW_H,
                        label, font_size=FS, bold=is_ms, color=lbl_color,
                        alignment=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)

    # ── Bars and milestones ───────────────────────────────────────
    for r, act in enumerate(activities):
        phase_key, label, bar_start, bar_end, is_ms, is_hdr, bar_clr_ov = act
        y = DATA_Y + r * ROW_H
        bar_color = bar_clr_ov or PC[phase_key]["bar"]

        if is_ms:
            # Diamond at milestone month
            mc = bar_start
            if col_start <= mc <= col_end:
                i = mc - col_start
                cx = CHART_X + i * COL_W + COL_W / 2
                cy = y + ROW_H / 2
                d = min(ROW_H * 0.72, COL_W * 0.60)
                add_shape(slide, MSO_SHAPE.DIAMOND,
                          cx - d / 2, cy - d / 2, d, d,
                          fill_color=MILESTONE_COLOR, line_color=None)
        else:
            vis_s = max(bar_start, col_start)
            vis_e = min(bar_end,   col_end)
            if vis_s > vis_e:
                continue
            i_s = vis_s - col_start
            i_e = vis_e - col_start
            bx = CHART_X + i_s * COL_W
            bw = (i_e - i_s + 1) * COL_W

            if is_hdr:
                # Phase header: full-height color band
                add_shape(slide, MSO_SHAPE.RECTANGLE, bx, y, bw, ROW_H,
                          fill_color=bar_color, line_color=None)
            else:
                bar_h = ROW_H - 2 * BAR_INSET
                bar_y = y + BAR_INSET
                add_shape(slide, MSO_SHAPE.RECTANGLE, bx, bar_y, bw, bar_h,
                          fill_color=bar_color, line_color=None)


# ──────────────────────────────────────────────────────────────────
# LEGEND  (drawn on overview slide below chart)
# ──────────────────────────────────────────────────────────────────

def draw_legend(slide, y_top=6.52):
    """Paint the phase legend at the given y position."""
    entries = [
        ("Phase 0 · Foundation",       "P0"),
        ("Phase 1 · Design",           "P1"),
        ("Phase 2 · Pilot (Wave 1)",   "P2"),
        ("Phase 3 · Scale — Wave 2",   "P3W2"),
        ("Phase 3 · Scale — Wave 3",   "P3W3"),
        ("Phase 4 · Sustain & Mature", "P4"),
        ("◆ Key Milestone",            None),
    ]
    LX    = 0.35
    TOTAL = 12.63
    item_w = TOTAL / 4
    item_h = 0.26

    for i, (label, pk) in enumerate(entries):
        col = i % 4
        row = i // 4
        ix = LX + col * item_w
        iy = y_top + row * item_h

        if pk is None:
            # Milestone diamond swatch
            d = 0.14
            cx = ix + 0.05 + d / 2
            cy = iy + item_h / 2
            add_shape(slide, MSO_SHAPE.DIAMOND,
                      cx - d / 2, cy - d / 2, d, d,
                      fill_color=MILESTONE_COLOR)
            add_textbox(slide, ix + 0.24, iy, item_w - 0.26, item_h,
                        label, font_size=7, bold=True, color=MILESTONE_COLOR)
        else:
            color = PC[pk]["bar"]
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      ix + 0.05, iy + 0.06, 0.15, 0.12,
                      fill_color=color)
            add_textbox(slide, ix + 0.24, iy, item_w - 0.26, item_h,
                        label, font_size=7, color="1D1D1D")


# ──────────────────────────────────────────────────────────────────
# MILESTONES SUMMARY TABLE  (drawn on a single summary slide)
# ──────────────────────────────────────────────────────────────────

def month_label(idx):
    """Return 'Mon YYYY' for a 0-based month index."""
    mo, yr = MONTHS[idx]
    return f"{mo} {yr}"


def create_milestones_slide(prs, page_number, logo_path):
    """Create a slide listing all key milestones."""
    from slide_builder import create_comparison_table
    milestones = [
        (a[1].lstrip("◆ ").strip(), month_label(a[2]), a[0])
        for a in ACTIVITIES if a[4]   # is_milestone
    ]
    rows = []
    for label, date, phase_key in milestones:
        phase_name = {
            "P0":   "Phase 0 · Foundation",
            "P1":   "Phase 1 · Design",
            "P2":   "Phase 2 · Pilot (Wave 1)",
            "P3W2": "Phase 3 · Scale — Wave 2",
            "P3W3": "Phase 3 · Scale — Wave 3",
            "P4":   "Phase 4 · Sustain & Mature",
        }.get(phase_key, phase_key)
        rows.append([date, label, phase_name])

    slide = blank_slide(prs)
    create_comparison_table(
        slide,
        title="Key Milestones — Agile Rollout Program",
        headers=["Target Date", "Milestone", "Phase"],
        rows=rows,
        col_widths=[1.4, 6.5, 4.7],
        logo_path=logo_path,
        page_number=page_number,
    )
    return slide


# ──────────────────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────────────────

def main():
    logo = find_infineon_logo()

    prs = create_presentation(
        title="Agile Scrum Rollout — Automotive Semiconductor R&D",
        author="Infineon Technologies",
    )

    # ── SLIDE 1: Title ───────────────────────────────────────────
    create_title_slide(
        prs,
        "Agile Scrum Rollout\nAutomotive Semiconductor R&D",
        subtitle="Program Gantt Chart  2026–2028  ·  1500 Engineers",
        logo_path=logo,
    )

    # ── SLIDE 2: Program Overview (phases only, all 28 months) ───
    phase_rows = [a for a in ACTIVITIES if a[5]]   # is_phase_hdr == True
    slide2 = blank_slide(prs)
    draw_gantt_slide(
        slide2,
        "Agile Transformation — Program Overview  (Apr 2026 – May 2028)",
        phase_rows, 0, 25,
        page_number=2, logo_path=logo,
        reserve_bottom=0.70,           # space for legend
    )
    draw_legend(slide2, y_top=6.52)

    # ── SLIDE 3: Phase 0 & 1 — Foundation & Design ───────────────
    p01 = [a for a in ACTIVITIES if a[0] in ("P0", "P1")]
    slide3 = blank_slide(prs)
    draw_gantt_slide(
        slide3,
        "Phase 0 & 1 — Foundation & Design  (Apr – Sep 2026)",
        p01, 0, 5,
        page_number=3, logo_path=logo,
    )

    # ── SLIDE 4: Phase 2 — Pilot Wave 1 ──────────────────────────
    p2 = [a for a in ACTIVITIES if a[0] == "P2"]
    slide4 = blank_slide(prs)
    draw_gantt_slide(
        slide4,
        "Phase 2 — Pilot Wave 1  (Jun – Dec 2026)",
        p2, 2, 8,
        page_number=4, logo_path=logo,
    )

    # ── SLIDE 5: Phase 3 · Wave 2 — Scale ────────────────────────
    p3w2 = [a for a in ACTIVITIES if a[0] == "P3W2"]
    slide5 = blank_slide(prs)
    draw_gantt_slide(
        slide5,
        "Phase 3 · Wave 2 — Scale  (Dec 2026 – Aug 2027)",
        p3w2, 8, 15,
        page_number=5, logo_path=logo,
    )

    # ── SLIDE 6: Phase 3 Wave 3 + Phase 4 ────────────────────────
    p3w3_p4 = [a for a in ACTIVITIES if a[0] in ("P3W3", "P4")]
    slide6 = blank_slide(prs)
    draw_gantt_slide(
        slide6,
        "Phase 3 · Wave 3 & Phase 4 — Scale & Sustain  (Jul 2027 – May 2028)",
        p3w3_p4, 15, 25,
        page_number=6, logo_path=logo,
    )

    # ── SLIDE 7: Key Milestones Summary ──────────────────────────
    create_milestones_slide(prs, page_number=7, logo_path=logo)

    # ── Save  (next to the source Excel file) ────────────────────
    script_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.normpath(
        os.path.join(script_dir, "..", "..", "..", "Agile_Rollout_Gantt_Presentation.pptx")
    )
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
